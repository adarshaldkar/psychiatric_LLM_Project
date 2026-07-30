"""PPTX parser — slide title = heading, bullet points = paragraph body."""
import logging
from typing import List

from app.rag.parsers.base_parser import BaseParser, ParsedDocument, ParsedPage

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    def can_parse(self, file_type: str) -> bool:
        return file_type.lower() == 'pptx'

    def parse(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise RuntimeError("python-pptx not installed")

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Cannot open PPTX: {e}")

        pages: List[ParsedPage] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Detect slide title shapes or bold large text as headings
                    is_title = (
                        shape.shape_type == 13 or                     # title placeholder
                        'title' in shape.name.lower() or
                        any(
                            run.font.bold and run.font.size and run.font.size >= Pt(20)
                            for run in para.runs
                            if run.font.size
                        )
                    )

                    if is_title:
                        parts.append(f'\n## {text}\n')
                    else:
                        parts.append(text)

            slide_text = '\n'.join(parts).strip()
            if slide_text:
                pages.append(ParsedPage(
                    page_number=slide_num,
                    text=slide_text,
                    used_ocr=False
                ))

        logger.info(f"PPTX parsed: {len(prs.slides)} slides → {len(pages)} non-empty — {file_path}")

        return ParsedDocument(pages=pages, file_type='pptx')
